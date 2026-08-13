# -*- coding: utf-8 -*-
"""生成《难得读书》6页竞赛 PPT，带排版配色 + 二维码。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- 主题色（与 App 薄荷绿一致）----
DEEP    = RGBColor(0x06, 0x46, 0x3C)   # 深墨绿（标题）
GREEN   = RGBColor(0x0B, 0x6B, 0x57)   # 主绿
MINT    = RGBColor(0x14, 0xB8, 0xA6)   # 薄荷
LIGHT   = RGBColor(0xE9, 0xFC, 0xF5)   # 浅背景
GREY    = RGBColor(0x3A, 0x4A, 0x47)   # 正文灰
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CARDBG  = RGBColor(0xF3, 0xFB, 0xF8)

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def bg(s, color=WHITE):
    r = rect(s, 0, 0, SW, SH, color)
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return r

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold)"""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.line_spacing = 1.15
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = FONT
    return tb

def sidebar(s):
    rect(s, 0, 0, Inches(0.28), SH, MINT)

def page_num(s, n):
    txt(s, SW-Inches(1.2), SH-Inches(0.55), Inches(0.9), Inches(0.4),
        [[(f"0{n} / 06", 11, GREEN, False)]], align=PP_ALIGN.RIGHT)

def bullet(size=15):
    return size

# ============ 第1页 定位（封面）============
s = slide(); bg(s, LIGHT)
rect(s, 0, 0, SW, Inches(2.5), DEEP)
txt(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(1.0),
    [[("难得读书", 44, WHITE, True)]])
txt(s, Inches(0.92), Inches(1.55), Inches(11.5), Inches(0.6),
    [[("选一本书，附身角色，感受风雪，照见自己", 18, MINT, True)]])

txt(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.7),
    [[("一句话定位：", 18, DEEP, True), ("我们让书成为一面照见自己的镜子。", 18, GREEN, True)]])

# 三层卡片
labels = [
    ("初级 · 摘要 / 问答", "帮你知道书讲了什么（赛道明确不鼓励）", RGBColor(0x9C,0xB5,0xAF)),
    ("中级 · 角色扮演 / 进入剧情", "帮你理解人物为什么这么做", RGBColor(0x4F,0x9E,0x8C)),
    ("高级（我们在这）· 用现代价值观拆解古人抉择", "让书成为照见自己的镜子", MINT),
]
cy = Inches(3.65)
for i, (t, d, c) in enumerate(labels):
    y = cy + Inches(i*0.92)
    rect(s, Inches(0.9), y, Inches(9.6), Inches(0.78), CARDBG, line=c)
    rect(s, Inches(0.9), y, Inches(0.14), Inches(0.78), c)
    txt(s, Inches(1.2), y+Inches(0.06), Inches(9.2), Inches(0.7),
        [[(t+"　", 15, DEEP, True), (d, 13, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)

txt(s, Inches(0.9), Inches(6.55), Inches(9.6), Inches(0.7),
    [[("我们做的不是“读一本书”，而是让书对你产生仅靠阅读得不到的、活的意义。", 14, GREEN, True)]])

# 二维码
qr = "/home/user/mybook/public/qr/mybook-qr.png"
if os.path.exists(qr):
    s.shapes.add_picture(qr, Inches(11.0), Inches(3.7), Inches(1.7), Inches(1.7))
    txt(s, Inches(10.8), Inches(5.45), Inches(2.1), Inches(0.6),
        [[("扫码即玩", 12, DEEP, True)], [("mybook-3eb25c73.eazo.dev", 8.5, GREEN, False)]],
        align=PP_ALIGN.CENTER)
page_num(s, 1)

# ============ 通用内容页 ============
def content_page(n, title, sub_hi, sub_lo, sections, closing):
    s = slide(); bg(s, WHITE); sidebar(s)
    txt(s, Inches(0.7), Inches(0.5), Inches(12), Inches(1.0),
        [[(title, 30, DEEP, True)]])
    rect(s, Inches(0.72), Inches(1.35), Inches(1.4), Inches(0.08), MINT)
    y = Inches(1.75)
    for (head, items) in sections:
        if head:
            txt(s, Inches(0.7), y, Inches(12), Inches(0.5),
                [[(head, 17, GREEN, True)]])
            y += Inches(0.55)
        for it in items:
            txt(s, Inches(1.0), y, Inches(11.4), Inches(0.55),
                [[("• ", 15, MINT, True), (it, 15, GREY, False)]])
            y += Inches(0.56)
        y += Inches(0.15)
    if closing:
        rect(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.75), LIGHT)
        txt(s, Inches(0.95), Inches(6.42), Inches(11.5), Inches(0.62),
            [[(closing, 15, DEEP, True)]], anchor=MSO_ANCHOR.MIDDLE)
    page_num(s, n)
    return s

# 第2页
content_page(2, "谁在“想读经典，却读不进去”？", None, None, [
    ("我们的用户", [
        "25–35 岁、有一定阅历的年轻人：想读经典、买了书，却总在第几页卡住",
        "在生活 / 职场里正面临真实抉择、渴望“照见自己”、想更懂自己的人",
        "反感说教式知识付费，想要有参与感、有共鸣的深度内容",
    ]),
    ("他们的痛点", [
        "名著=“该读但读不下去”：距离感强、说教味重、缺乏参与感",
        "传统方式（讲解、速读、读后感）本质是“别人嚼过喂给你”，你始终是旁观者",
    ]),
], "名著离我们太远，因为我们从来没“成为”过书里的人。")

# 第3页
content_page(3, "让你“附身”角色，替他做每一次艰难抉择", None, None, [
    (None, [
        "选一本书 → 附身一个角色（林冲、祥子、诸葛亮、安娜……）",
        "在忠于原著的关键情节里，由你亲自做出每一次两难抉择",
        "情节锚点忠于原著（隆中对→托孤→北伐的顺序不会乱），但内心与选择由你长出",
        "走到结尾，你的每一次选择汇成一份“照见你自己”的报告",
    ]),
], "核心机制：原著的骨架 + 你的血肉 = 一次借角色照见自己的旅程。")

# 第4页 价值观
content_page(4, "不是娱乐游戏，是一面镜子", None, None, [
    (None, [
        "原著诚实底线：只认原著小说，拒绝影视 / 游戏改编污染（不会冒出“紫霞仙子”）；不认识的书宁可诚实说“暂无法可靠还原”",
        "真实的自我暴露：全程不剧透“在测什么”，你才会做出真实选择，结局才照见真实的你",
        "意义感驱动留存：趣味性之外，给用户“更看清自己”的深层动机——区别于快消游戏的护城河",
    ]),
], "别人做的是“内容消费”，我们做的是“自我觉察”。")

# 第5页 亮点
content_page(5, "一个已经跑通的完整体验", None, None, [
    (None, [
        "覆盖数十部中外经典，每本多个可选角色，支持自填角色",
        "因果连贯的剧情推演，拒绝凭空转折",
        "“极端选择付出代价但不中断”，人人都能走到属于自己的结局报告",
        "结局镜像报告：把角色一生的命题，轻轻折射回你自己的现实",
        "移动端优先，扫码即玩，可分享故事与结果",
    ]),
], "建议配图：手机端截图 —— 选书页 / 剧情页 / 结果报告页")

# ============ 第6页 结尾 ============
s = slide(); bg(s, DEEP)
txt(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
    [[("难得读书，难得照见自己", 36, WHITE, True)]])
rect(s, Inches(0.92), Inches(2.0), Inches(1.6), Inches(0.09), MINT)
txt(s, Inches(0.9), Inches(2.4), Inches(8.5), Inches(1.5),
    [[("• ", 16, MINT, True), ("愿景：让每一本经典，都成为一次照见自己的机会", 16, LIGHT, False)],
     [("• ", 16, MINT, True), ("让“读不进去的名著”，变成“停不下来的自己”", 16, LIGHT, False)]], sp_after=12)
txt(s, Inches(0.9), Inches(4.3), Inches(8.5), Inches(1.6),
    [[("现在就扫码，附身一个角色，走一遍——", 18, MINT, True)],
     [("看看最后照见的那个人，是不是你。", 18, MINT, True)]], sp_after=8)
if os.path.exists(qr):
    s.shapes.add_picture(qr, Inches(10.2), Inches(2.7), Inches(2.2), Inches(2.2))
    txt(s, Inches(10.0), Inches(4.95), Inches(2.6), Inches(0.6),
        [[("mybook-3eb25c73.eazo.dev", 11, WHITE, True)]], align=PP_ALIGN.CENTER)

out = "/home/user/mybook/public/难得读书-竞赛PPT.pptx"
prs.save(out)
print("SAVED", out, os.path.getsize(out), "bytes")
